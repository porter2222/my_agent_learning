"""
文档/图片处理流水线 - 多模态摘要、OCR、向量化、入库
支持 PDF、DOCX、PPTX、DOC 及多种图片格式
"""
import os
import io
import base64
import json
import tempfile
import shutil
from pathlib import Path
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Dict, Any

import fitz  # PyMuPDF
from openai import OpenAI
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

from docx2pdf import convert
from llm_factory import LLMClient
from config import KIMI_API_KEY

# Kimi OCR 客户端（月之暗面 API）
kimi_client = OpenAI(
    api_key=KIMI_API_KEY,
    base_url="https://api.moonshot.cn/v1",
)


def kimi_file_upload(file_path: str) -> str:
    """使用 Kimi 文件提取 API 进行 OCR"""
    try:
        file_object = kimi_client.files.create(
            file=Path(file_path), purpose="file-extract"
        )
        content = kimi_client.files.content(file_id=file_object.id).text
        kimi_client.files.delete(file_id=file_object.id)
        return json.loads(content).get("content", "")
    except Exception as e:
        print(f"[WARN] Kimi OCR 失败: {e}")
        return ""


class FileProcessor:
    """文档处理：PDF/DOCX/PPT 转图片 → 摘要 → OCR → 向量化 → 写入 Qdrant"""

    def __init__(self, qdrant_manager):
        self.ai_models = LLMClient(provider="qwen-cn", model="qwen-vl-max")
        self.qwen_embedding = LLMClient(provider="qwen-cn", model="text-embedding-v4")
        self.qdrant_manager = qdrant_manager
        self.executor = ThreadPoolExecutor(max_workers=10)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        self.temp_dir = Path(__file__).parent / f"temp_images_{timestamp}"
        self.temp_dir.mkdir(exist_ok=True)

    def save_image(self, image_data: bytes, filename: str, doc_name: str) -> str:
        """保存图片到临时目录，返回本地路径"""
        safe_doc_name = "".join(
            c for c in doc_name if c.isalnum() or c in (" ", "-", "_")
        ).rstrip()
        doc_dir = self.temp_dir / safe_doc_name
        doc_dir.mkdir(exist_ok=True, parents=True)
        image_path = doc_dir / filename
        image_path.write_bytes(image_data)
        return str(image_path)

    def convert_ppt_to_images(self, ppt_path: str) -> List[Image.Image]:
        """将 PPT 每页渲染为图片"""
        prs = Presentation(ppt_path)
        images = []
        width, height = 1920, 1080
        try:
            font = ImageFont.truetype("C:/Windows/Fonts/simhei.ttf", 24)
        except Exception:
            font = ImageFont.load_default()

        for slide in prs.slides:
            img = Image.new("RGB", (width, height), (255, 255, 255))
            draw = ImageDraw.Draw(img)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    left = int(shape.left * width / 9144000)
                    top = int(shape.top * height / 9144000)
                    draw.text((left, top), shape.text, font=font, fill=(0, 0, 0))
                if getattr(shape, "shape_type", None) == 13:
                    try:
                        blob = io.BytesIO(shape.image.blob)
                        si = Image.open(blob).resize(
                            (
                                int(shape.width * width / 9144000),
                                int(shape.height * height / 9144000),
                            )
                        )
                        img.paste(
                            si,
                            (
                                int(shape.left * width / 9144000),
                                int(shape.top * height / 9144000),
                            ),
                        )
                    except Exception:
                        pass
            images.append(img)
        return images

    def convert_to_pdf(self, file_content: bytes, filename: str) -> bytes:
        """DOCX/PPTX 转 PDF，PDF 直接返回"""
        ext = os.path.splitext(filename)[1].lower()
        if ext == ".pdf":
            return file_content
        temp_in = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
        temp_out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        try:
            temp_in.write(file_content)
            temp_in.close()
            if ext == ".docx":
                convert(temp_in.name, temp_out.name)
            elif ext in (".pptx", ".ppt"):
                images = self.convert_ppt_to_images(temp_in.name)
                if images:
                    images[0].save(
                        temp_out.name,
                        save_all=True,
                        append_images=images[1:],
                        resolution=100.0,
                    )
            else:
                raise ValueError(f"不支持格式: {ext}")
            return Path(temp_out.name).read_bytes()
        finally:
            if os.path.exists(temp_in.name):
                os.remove(temp_in.name)
            if os.path.exists(temp_out.name):
                os.remove(temp_out.name)

    def pdf_to_images(self, pdf_content: bytes) -> List[bytes]:
        """PDF 每页转 PNG 字节"""
        tmp = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        tmp.write(pdf_content)
        tmp.close()
        try:
            doc = fitz.open(tmp.name)
            images = [
                doc[page].get_pixmap(matrix=fitz.Matrix(2, 2)).tobytes("png")
                for page in range(len(doc))
            ]
            doc.close()
            return images
        finally:
            os.remove(tmp.name)

    def process_single_image(
        self,
        image_data: bytes,
        image_index: int,
        doc_name: str,
        filename: str,
    ) -> Dict[str, Any]:
        """单页：摘要 → OCR → 向量化 → 写入 Qdrant"""
        image_filename = f"page_{image_index + 1}.png"
        image_path = self.save_image(image_data, image_filename, doc_name)
        base64_image = base64.b64encode(image_data).decode("utf-8")
        summary_text = self.ai_models.qwen_vision(
            base64_image,
            "请对图片做摘要，提取核心内容，200 字以内，直接输出摘要。",
        )
        tmp_img = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        tmp_img.write(image_data)
        tmp_img.close()
        try:
            origin_text = kimi_file_upload(tmp_img.name)
        finally:
            os.remove(tmp_img.name)
        vector = self.qwen_embedding.embedding(summary_text)
        metadata = {
            "image_path": image_path,
            "summary_text": summary_text,
            "origin_text": origin_text,
            "image_index": image_index,
        }
        self.qdrant_manager.store_vectors(vector=vector, metadata=metadata)
        return metadata

    def process_file_content(self, file_content: bytes, filename: str) -> dict:
        """文档整体处理：转 PDF → 转图 → 多线程处理每页"""
        doc_name = os.path.splitext(filename)[0]
        pdf_content = self.convert_to_pdf(file_content, filename)
        images = self.pdf_to_images(pdf_content)
        futures = [
            self.executor.submit(
                self.process_single_image, img, i, doc_name, filename
            )
            for i, img in enumerate(images)
        ]
        results = []
        for f in as_completed(futures):
            try:
                results.append(f.result())
            except Exception as e:
                print(f"[WARN] 处理失败: {e}")
        results.sort(key=lambda x: x["image_index"])
        return {
            "image_paths": [r["image_path"] for r in results],
            "summaries": [r["summary_text"] for r in results],
            "original_texts": [r["origin_text"] for r in results],
        }

    def process_image_file(self, image_content: bytes, filename: str) -> dict:
        """单张图片处理"""
        doc_name = os.path.splitext(filename)[0]
        r = self.process_single_image(image_content, 0, doc_name, filename)
        return {
            "image_paths": [r["image_path"]],
            "summaries": [r["summary_text"]],
            "original_texts": [r["origin_text"]],
        }

    def cleanup(self):
        """清理临时目录"""
        if self.temp_dir.exists():
            shutil.rmtree(self.temp_dir)
